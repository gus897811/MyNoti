"""MyNoti 발표용 PPT 생성."""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap, qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

ROOT = Path(__file__).resolve().parent
ASSETS = ROOT / "assets"
OUT = ROOT / "MyNoti_Presentation.pptx"

W = Inches(13.333)
H = Inches(7.5)
FONT = "Malgun Gothic"

NAVY = RGBColor(0x0F, 0x17, 0x2A)
NAVY2 = RGBColor(0x1E, 0x29, 0x3B)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
BLUE_DARK = RGBColor(0x25, 0x63, 0xEB)
BLUE_SOFT = RGBColor(0xDB, 0xEA, 0xFE)
BG = RGBColor(0xF7, 0xF8, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT = RGBColor(0x11, 0x18, 0x27)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LINE = RGBColor(0xE8, 0xEA, 0xF0)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
ORANGE_BG = RGBColor(0xFF, 0xF8, 0xF1)
GREEN = RGBColor(0x05, 0x96, 0x69)
GREEN_BG = RGBColor(0xEC, 0xFD, 0xF5)
RED = RGBColor(0xDC, 0x26, 0x26)
RED_BG = RGBColor(0xFE, 0xF2, 0xF2)
KAKAO = RGBColor(0xFE, 0xE5, 0x00)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)


def rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def set_fill(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def set_line(shape, color: RGBColor, pt: float = 1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(pt)


def roundness(shape, value: float = 0.12):
    try:
        shape.adjustments[0] = value
    except Exception:
        pass


def add_rect(slide, x, y, w, h, color, radius=None):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(kind, x, y, w, h)
    set_fill(s, color)
    if radius is not None:
        roundness(s, radius)
    return s


def add_outline(slide, x, y, w, h, fill, line, radius=0.1, lw=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    set_line(s, line, lw)
    roundness(s, radius)
    return s


def _set_run(run, size, color, bold=False, font=FONT):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font)


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=16,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font=FONT,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    try:
        box.text_frame._txBody.bodyPr.set("anchor", {MSO_ANCHOR.TOP: "t", MSO_ANCHOR.MIDDLE: "ctr", MSO_ANCHOR.BOTTOM: "b"}[anchor])
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    _set_run(run, size, color, bold, font)
    return box


def add_lines(slide, x, y, w, h, lines, size=15, color=TEXT, bold=False, align=PP_ALIGN.LEFT, spacing=1.08):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        p.line_spacing = spacing
        run = p.add_run()
        run.text = line
        _set_run(run, size, color, bold)
    return box


def content_bg(slide):
    add_rect(slide, 0, 0, W, H, BG)
    add_rect(slide, 0, 0, W, Inches(0.08), BLUE)


def footer(slide, page: int, total: int = 15):
    add_text(slide, Inches(0.55), Inches(7.18), Inches(4), Inches(0.22), "MyNoti", 11, MUTED, False)
    add_text(
        slide,
        Inches(11.4),
        Inches(7.18),
        Inches(1.4),
        Inches(0.22),
        f"{page}  /  {total}",
        11,
        MUTED,
        False,
        PP_ALIGN.RIGHT,
    )


def kicker(slide, text: str):
    add_text(slide, Inches(0.6), Inches(0.28), Inches(8), Inches(0.28), text, 12, BLUE, True)


def title(slide, text: str, y=0.52):
    add_text(slide, Inches(0.6), Inches(y), Inches(12.1), Inches(0.55), text, 28, TEXT, True)


def round_phone(src: Path, dst: Path, radius: int = 52):
    im = Image.open(src).convert("RGBA")
    w, h = im.size
    mask = Image.new("L", (w, h), 0)
    ImageDraw.Draw(mask).rounded_rectangle((0, 0, w - 1, h - 1), radius=radius, fill=255)
    im.putalpha(mask)
    # thin frame
    framed = Image.new("RGBA", (w + 16, h + 16), (0, 0, 0, 0))
    bg = Image.new("RGBA", (w + 16, h + 16), (0, 0, 0, 0))
    ImageDraw.Draw(bg).rounded_rectangle((0, 0, w + 15, h + 15), radius=radius + 8, fill=(17, 24, 39, 255))
    framed.alpha_composite(bg)
    framed.alpha_composite(im, (8, 8))
    framed.save(dst)


def phone(slide, src: Path, x, y, h):
    im = Image.open(src)
    ratio = im.width / im.height
    height = h
    width = height * ratio
    # shadow
    add_rect(slide, x + Inches(0.06), y + Inches(0.07), width, height, RGBColor(0xD1, 0xD5, 0xDB), 0.12)
    slide.shapes.add_picture(str(src), x, y, width, height)
    return width, height


def arrow_right(slide, x, y, w=Inches(0.28), h=Inches(0.18), color=BLUE):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    set_fill(s, color)
    return s


def chevron_box(slide, x, y, w, h, label, sub, fill, text_color=WHITE):
    s = add_rect(slide, x, y, w, h, fill, 0.14)
    add_text(slide, x + Inches(0.14), y + Inches(0.12), w - Inches(0.28), Inches(0.28), label, 13, text_color, True)
    add_text(slide, x + Inches(0.14), y + Inches(0.40), w - Inches(0.28), Inches(0.42), sub, 11, text_color if text_color != WHITE else RGBColor(0xDB, 0xEA, 0xFE))
    return s


# ── slides ──────────────────────────────────────────────────────────────────


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, 0, Inches(0.18), H, BLUE)
    # decorative blobs
    blob = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(-1.2), Inches(4.4), Inches(4.4))
    set_fill(blob, RGBColor(0x1E, 0x3A, 0x8A))
    blob2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.4), Inches(5.2), Inches(3.2), Inches(3.2))
    set_fill(blob2, RGBColor(0x17, 0x24, 0x54))

    add_text(s, Inches(0.9), Inches(1.55), Inches(10), Inches(0.35), "PROJECT  PRESENTATION", 13, BLUE, True)
    add_text(s, Inches(0.9), Inches(2.05), Inches(11), Inches(1.1), "MyNoti", 66, WHITE, True)
    add_text(s, Inches(0.9), Inches(3.2), Inches(10), Inches(0.45), "Smarter Alerts, Zero Noise.", 24, RGBColor(0x93, 0xC5, 0xFD), False)
    add_text(
        s,
        Inches(0.9),
        Inches(4.0),
        Inches(9.5),
        Inches(0.7),
        "잠에서 깨도, 중요한 알림만 남아 있게.\n대학생을 위한 스마트 알림 정리 앱",
        16,
        RGBColor(0xCB, 0xD5, 0xE1),
        False,
    )
    add_rect(s, Inches(0.9), Inches(5.55), Inches(2.1), Inches(0.06), BLUE, 0.5)
    add_text(s, Inches(0.9), Inches(5.8), Inches(8), Inches(0.3), "알림 분석  ·  리마인더  ·  캘린더", 14, RGBColor(0x94, 0xA3, 0xB8))


def slide_agenda(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    kicker(s, "CONTENTS")
    title(s, "오늘 이야기할 것")

    items = [
        ("01", "기획 배경", "알림은 끄기도, 켜두기도 어려운 이유"),
        ("02", "주요 기능", "분석 · 알람 · 캘린더"),
        ("03", "시스템 구조", "알림이 들어와서 화면에 보이기까지"),
        ("04", "시연", "실제 앱으로 이어서 보여드립니다"),
    ]
    for i, (num, head, sub) in enumerate(items):
        x = Inches(0.6) + (i % 2) * Inches(6.25)
        y = Inches(1.45) + (i // 2) * Inches(2.45)
        add_outline(s, x, y, Inches(5.95), Inches(2.15), WHITE, LINE, 0.08, 1.25)
        add_text(s, x + Inches(0.35), y + Inches(0.35), Inches(2), Inches(0.45), num, 28, BLUE, True)
        add_text(s, x + Inches(0.35), y + Inches(0.95), Inches(5.2), Inches(0.4), head, 22, TEXT, True)
        add_text(s, x + Inches(0.35), y + Inches(1.4), Inches(5.2), Inches(0.4), sub, 14, MUTED)
    footer(s, 2)


def slide_morning(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    kicker(s, "01  기획 배경  ·  어느 아침")
    title(s, "눈 떴는데, 알림이 47개.")

    add_text(
        s,
        Inches(0.6),
        Inches(1.15),
        Inches(6.4),
        Inches(0.55),
        "잠자는 동안 카톡, 수업 앱, 카드 앱이\n한꺼번에 쌓여 있다.",
        16,
        MUTED,
    )

    notifs = [
        (KAKAO, NAVY, "카카오톡  ·  김경래", "ㅋㅋㅋㅋ ㄱㄱ", "방금"),
        (KAKAO, NAVY, "카카오톡  ·  광고", "[봄맞이] 세일 50% 오늘 마감!", "1분"),
        (BLUE, WHITE, "LearningX", "로그인되었습니다", "3분"),
        (NAVY2, WHITE, "신한카드", "[이벤트] 이번 주 캐시백 혜택", "8분"),
        (KAKAO, NAVY, "카카오톡  ·  팀장", "오늘 21시까지 GitHub 링크 제출…", "12분"),
        (BLUE, WHITE, "LearningX", "캡스톤 중간보고서가 등록되었습니다", "20분"),
    ]
    y0 = Inches(1.85)
    for i, (dot, _tc, app, body, when) in enumerate(notifs):
        y = y0 + Inches(i * 0.78)
        card = add_outline(s, Inches(0.6), y, Inches(6.55), Inches(0.70), WHITE, LINE, 0.16, 1.0)
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.78), y + Inches(0.20), Inches(0.30), Inches(0.30))
        set_fill(badge, dot)
        add_text(s, Inches(1.22), y + Inches(0.08), Inches(4.4), Inches(0.26), app, 11, MUTED, True)
        add_text(s, Inches(1.22), y + Inches(0.34), Inches(4.7), Inches(0.28), body, 14, TEXT, True)
        add_text(s, Inches(5.85), y + Inches(0.22), Inches(1.1), Inches(0.28), when, 11, MUTED, False, PP_ALIGN.RIGHT)

    # right drama card
    add_rect(s, Inches(7.55), Inches(1.85), Inches(5.15), Inches(4.75), NAVY, 0.08)
    add_text(s, Inches(7.9), Inches(2.15), Inches(4.5), Inches(0.3), "잠금화면 속마음", 12, BLUE, True)
    add_text(
        s,
        Inches(7.9),
        Inches(2.6),
        Inches(4.5),
        Inches(1.6),
        "“…그래서\n뭐가 중요한 건데?”",
        26,
        WHITE,
        True,
    )
    add_text(
        s,
        Inches(7.9),
        Inches(4.4),
        Inches(4.5),
        Inches(1.7),
        "쓸데없는 내용도 많고,\n처음부터 끝까지 읽기엔 너무 많다.\n\n중요한 건 분명 섞여 있는데\n찾으려면 전부 열어봐야 한다.",
        15,
        RGBColor(0xCB, 0xD5, 0xE1),
    )
    footer(s, 3)


def slide_dilemma(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    kicker(s, "01  기획 배경  ·  딜레마")
    title(s, "그렇다고 그 앱 알림을 끄기도 뭐하다")

    # left: off
    add_outline(s, Inches(0.6), Inches(1.4), Inches(5.95), Inches(4.55), WHITE, LINE, 0.07, 1.25)
    add_rect(s, Inches(0.6), Inches(1.4), Inches(0.14), Inches(4.55), RED)
    add_text(s, Inches(1.05), Inches(1.65), Inches(5.1), Inches(0.3), "알림을 끄면", 13, RED, True)
    add_text(s, Inches(1.05), Inches(2.05), Inches(5.1), Inches(0.7), "중요한 메시지를\n놓친다.", 26, TEXT, True)
    add_outline(s, Inches(1.05), Inches(3.05), Inches(5.05), Inches(2.4), RED_BG, RGBColor(0xFE, 0xCD, 0xD3), 0.1)
    add_text(s, Inches(1.3), Inches(3.25), Inches(4.6), Inches(0.28), "다음 날, 팀장에게", 12, RED, True)
    add_text(
        s,
        Inches(1.3),
        Inches(3.6),
        Inches(4.6),
        Inches(1.5),
        "“야 너 카톡 봤어?\n오늘 21시까지라니까??\n공지 안 읽었어?”",
        18,
        TEXT,
        True,
    )

    # right: on
    add_outline(s, Inches(6.8), Inches(1.4), Inches(5.95), Inches(4.55), WHITE, LINE, 0.07, 1.25)
    add_rect(s, Inches(6.8), Inches(1.4), Inches(0.14), Inches(4.55), ORANGE)
    add_text(s, Inches(7.25), Inches(1.65), Inches(5.1), Inches(0.3), "알림을 켜두면", 13, ORANGE, True)
    add_text(s, Inches(7.25), Inches(2.05), Inches(5.1), Inches(0.7), "정리가\n안 된다.", 26, TEXT, True)
    add_outline(s, Inches(7.25), Inches(3.05), Inches(5.05), Inches(2.4), ORANGE_BG, RGBColor(0xFE, 0xD7, 0xAA), 0.1)
    add_text(s, Inches(7.5), Inches(3.25), Inches(4.6), Inches(0.28), "한 줄에 섞여 있는 것들", 12, ORANGE, True)
    add_lines(
        s,
        Inches(7.5),
        Inches(3.65),
        Inches(4.6),
        Inches(1.6),
        ["광고 · 잡담 · 로그인 알림", "과제 마감 · 수업 공지 · 정산", "중요한 것과 아닌 것의 경계가 없다"],
        16,
        TEXT,
    )
    footer(s, 4)


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    kicker(s, "01  기획 배경  ·  한 줄")
    title(s, "끄면 놓치고, 켜면 묻힌다.")

    points = [
        ("쌓인다", "자고 일어나면 알림이 한꺼번에 올라온다."),
        ("섞인다", "광고 · 잡담 · 과제가 같은 목록에 있다."),
        ("못 끈다", "그 앱이 중요한 메시지를 보낼 수 있다."),
    ]
    for i, (h, b) in enumerate(points):
        x = Inches(0.6) + i * Inches(4.15)
        add_outline(s, x, Inches(1.45), Inches(3.95), Inches(2.15), WHITE, LINE, 0.08)
        add_text(s, x + Inches(0.3), Inches(1.7), Inches(1.2), Inches(0.35), f"0{i+1}", 16, BLUE, True)
        add_text(s, x + Inches(0.3), Inches(2.15), Inches(3.35), Inches(0.45), h, 24, TEXT, True)
        add_text(s, x + Inches(0.3), Inches(2.7), Inches(3.35), Inches(0.6), b, 14, MUTED)

    add_rect(s, Inches(0.6), Inches(3.95), Inches(12.1), Inches(2.55), NAVY, 0.07)
    add_text(s, Inches(0.95), Inches(4.25), Inches(11.4), Inches(0.3), "그래서 만든 것", 13, BLUE, True)
    add_text(s, Inches(0.95), Inches(4.65), Inches(11.4), Inches(0.7), "알림은 그대로 받되, 중요한 것만 정리해서 보여준다.", 24, WHITE, True)
    add_text(
        s,
        Inches(0.95),
        Inches(5.45),
        Inches(11.4),
        Inches(0.6),
        "MyNoti는 지정한 앱의 알림을 모아 LLM으로 분류·요약하고, 과제·마감·수업만 홈 · 요약 · 캘린더 · 리마인더로 남깁니다.",
        15,
        RGBColor(0xCB, 0xD5, 0xE1),
    )
    footer(s, 5)


def slide_features(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    kicker(s, "02  주요 기능")
    title(s, "세 가지로 알림을 다시 씁니다")

    cards = [
        ("2.1", "알림 받아서 분석", "원문을 모아 요약하고\n중요 · 과제 · 소통으로 분류합니다.\n잡담과 광고는 걸러냅니다.", BLUE),
        ("2.2", "알람 설정", "다시 알려줄 시각을 걸고,\n중요/숨김 키워드와\n감시할 앱을 직접 고릅니다.", ORANGE),
        ("2.3", "캘린더", "알림에서 뽑은 마감이\n달력에 올라갑니다.\n직접 일정을 추가할 수도 있습니다.", GREEN),
    ]
    for i, (num, head, body, accent) in enumerate(cards):
        x = Inches(0.6) + i * Inches(4.15)
        add_outline(s, x, Inches(1.4), Inches(3.95), Inches(5.05), WHITE, LINE, 0.07)
        add_rect(s, x, Inches(1.4), Inches(3.95), Inches(0.12), accent)
        add_text(s, x + Inches(0.3), Inches(1.8), Inches(3.3), Inches(0.35), num, 14, accent, True)
        add_text(s, x + Inches(0.3), Inches(2.25), Inches(3.35), Inches(1.0), head, 24, TEXT, True)
        add_text(s, x + Inches(0.3), Inches(3.45), Inches(3.35), Inches(2.2), body, 16, MUTED)
    footer(s, 6)


def slide_analyze(prs, home_img: Path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    kicker(s, "02  주요 기능  ·  2.1")
    title(s, "알림을 받아서, 분석합니다")

    bullets = [
        ("수집", "카톡 · LearningX · 카드앱 등\n지정한 앱의 알림만 모읍니다."),
        ("1차 필터", "잡담 · 광고 · 로그인 알림처럼\n볼 필요 없는 것은 버립니다."),
        ("본 분석", "제목 요약, 유형, 중요도,\n마감 시각을 뽑아냅니다."),
        ("표시", "홈에서 카드로 보고,\n중요 알림은 한눈에 띄게 합니다."),
    ]

    for i, (h, b) in enumerate(bullets):
        y = Inches(1.35) + i * Inches(1.28)
        add_outline(s, Inches(0.6), y, Inches(7.35), Inches(1.16), WHITE, LINE, 0.14)
        n = add_rect(s, Inches(0.82), y + Inches(0.32), Inches(0.52), Inches(0.52), BLUE_SOFT, 0.2)
        add_text(s, Inches(0.82), y + Inches(0.40), Inches(0.52), Inches(0.38), f"{i+1}", 14, BLUE, True, PP_ALIGN.CENTER)
        add_text(s, Inches(1.55), y + Inches(0.18), Inches(6.1), Inches(0.32), h, 16, TEXT, True)
        add_text(s, Inches(1.55), y + Inches(0.52), Inches(6.1), Inches(0.52), b, 13, MUTED)

    phone(s, home_img, Inches(8.35), Inches(1.28), Inches(5.55))
    footer(s, 7)


def slide_summary(prs, summary_img: Path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    kicker(s, "02  주요 기능  ·  2.1")
    title(s, "분석이 끝나면, 오늘이 한 장으로 남습니다")

    phone(s, summary_img, Inches(0.7), Inches(1.28), Inches(5.55))

    items = [
        ("숫자로 먼저", "중요한 알림 · 과제 · 예정된 일정을\n세 장 카드로 보여줍니다."),
        ("가장 급한 일", "오늘 처리해야 할 일을\n마감 가까운 순으로 골라줍니다."),
        ("AI Insight", "지금 무엇을 먼저 할지\n한 문장으로 추천합니다."),
        ("다시 알림", "직접 걸어 둔 리마인더를\n시간대별로 모아 보여줍니다."),
    ]
    for i, (h, b) in enumerate(items):
        y = Inches(1.35) + i * Inches(1.28)
        add_outline(s, Inches(3.85), y, Inches(8.85), Inches(1.16), WHITE, LINE, 0.12)
        add_text(s, Inches(4.15), y + Inches(0.18), Inches(8.2), Inches(0.32), h, 16, TEXT, True)
        add_text(s, Inches(4.15), y + Inches(0.52), Inches(8.2), Inches(0.52), b, 14, MUTED)
    footer(s, 8)


def slide_alarm(prs, kw_img: Path, apps_img: Path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    kicker(s, "02  주요 기능  ·  2.2")
    title(s, "다시 알려주고, 규칙을 내가 정합니다")

    # three feature cards on left
    blocks = [
        (ORANGE, "리마인더", "알림 하나에 ‘이때 다시 울려줘’를 겁니다.\n시간이 되면 MyNoti가 다시 알려줍니다."),
        (BLUE, "중요 키워드", "과제 · 시험 · 마감 · 긴급 · 장학금.\n이 단어가 있으면 중요으로 표시합니다."),
        (MUTED, "숨길 키워드", "광고 · 세일 · 프로모션 · spam.\n목록에서 숨깁니다. 원문은 기기에 남습니다."),
        (GREEN, "타겟 앱", "감시할 앱을 켜고 끕니다.\n꺼 둔 앱의 새 알림은 모으지 않습니다."),
    ]
    for i, (c, h, b) in enumerate(blocks):
        col, row = i % 2, i // 2
        x = Inches(0.6) + col * Inches(4.55)
        y = Inches(1.32) + row * Inches(2.55)
        add_outline(s, x, y, Inches(4.35), Inches(2.35), WHITE, LINE, 0.1)
        add_rect(s, x + Inches(0.28), y + Inches(0.28), Inches(0.18), Inches(0.18), c, 0.5)
        add_text(s, x + Inches(0.58), y + Inches(0.22), Inches(3.5), Inches(0.32), h, 16, TEXT, True)
        add_text(s, x + Inches(0.28), y + Inches(0.7), Inches(3.8), Inches(1.35), b, 14, MUTED)

    phone(s, kw_img, Inches(9.85), Inches(1.28), Inches(5.55))
    footer(s, 9)


def slide_calendar(prs, cal_img: Path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    kicker(s, "02  주요 기능  ·  2.3")
    title(s, "알림 속 마감이, 달력으로 올라갑니다")

    add_text(
        s,
        Inches(0.6),
        Inches(1.2),
        Inches(7.4),
        Inches(0.55),
        "분석이 날짜를 찾으면 캘린더에 자동으로 찍힙니다.\n직접 ‘+ 일정 추가’로 넣을 수도 있습니다.",
        15,
        MUTED,
    )

    rows = [
        ("자동 일정", "LLM이 알림에서 뽑은 마감 시각.\n예: “오늘 Basic반 4시 수업”"),
        ("수동 일정", "사용자가 직접 제목과 시간을 넣습니다."),
        ("하루 보기", "날짜를 누르면 그날 알림·일정이\n카드로 펼쳐집니다."),
        ("색 점", "날짜 아래 점으로 일정이 있는 날을\n한눈에 봅니다."),
    ]
    for i, (h, b) in enumerate(rows):
        y = Inches(1.95) + i * Inches(1.15)
        add_outline(s, Inches(0.6), y, Inches(7.45), Inches(1.05), WHITE, LINE, 0.14)
        add_text(s, Inches(0.9), y + Inches(0.16), Inches(6.9), Inches(0.3), h, 16, TEXT, True)
        add_text(s, Inches(0.9), y + Inches(0.48), Inches(6.9), Inches(0.48), b, 13, MUTED)

    phone(s, cal_img, Inches(8.5), Inches(1.28), Inches(5.55))
    footer(s, 10)


def slide_architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    kicker(s, "03  시스템 구조")
    title(s, "역할이 갈린 세 층")

    layers = [
        (NAVY, WHITE, "사용자 폰", "MyNoti Android", "알림을 가로채고, Room에 원문을 보관하고,\n홈 · 요약 · 캘린더 · 리마인더로 보여줍니다."),
        (BLUE, WHITE, "우리 서버", "MyNoti Backend", "FastAPI 게이트웨이.\n분석만 하고, 알림을 저장하지 않습니다."),
        (RGBColor(0x10, 0xB9, 0x81), WHITE, "LLM", "OpenAI", "1차: 보여줄 가치 있나?\n2차: 요약 · 유형 · 중요도 · 마감."),
    ]
    for i, (fill, tc, k, h, b) in enumerate(layers):
        x = Inches(0.6) + i * Inches(4.15)
        add_rect(s, x, Inches(1.4), Inches(3.95), Inches(3.35), fill, 0.07)
        add_text(s, x + Inches(0.3), Inches(1.65), Inches(3.35), Inches(0.3), k, 13, RGBColor(0xBF, 0xDB, 0xFE) if fill != RGBColor(0x10, 0xB9, 0x81) else RGBColor(0xD1, 0xFA, 0xE5), True)
        add_text(s, x + Inches(0.3), Inches(2.05), Inches(3.35), Inches(0.7), h, 22, WHITE, True)
        add_text(s, x + Inches(0.3), Inches(2.9), Inches(3.35), Inches(1.5), b, 14, RGBColor(0xE2, 0xE8, 0xF0))

    notes = [
        ("원문은 폰에만", "사용자 알림의 원본은 Room(mynoti.db)입니다."),
        ("키도 갈린다", "앱은 API 키만, OpenAI 키는 서버에만 있습니다."),
        ("키워드는 앱에서", "Highlight / Mute는 서버가 모릅니다. 목록 표시 때 적용합니다."),
    ]
    for i, (h, b) in enumerate(notes):
        x = Inches(0.6) + i * Inches(4.15)
        add_outline(s, x, Inches(5.0), Inches(3.95), Inches(1.55), WHITE, LINE, 0.12)
        add_text(s, x + Inches(0.25), Inches(5.18), Inches(3.45), Inches(0.32), h, 14, TEXT, True)
        add_text(s, x + Inches(0.25), Inches(5.52), Inches(3.45), Inches(0.8), b, 13, MUTED)
    footer(s, 11)


def slide_flow(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    kicker(s, "03  시스템 구조  ·  흐름")
    title(s, "알림이 화면에 보이기까지")

    steps = [
        ("1", "발생", "카톡 · 수업앱\n카드앱 알림"),
        ("2", "수집", "Listener가\n원문을 가로챔"),
        ("3", "저장", "Room에 PENDING\n으로 넣음"),
        ("4", "분석", "Worker가\n백엔드로 요청"),
        ("5", "판단", "LLM 필터\n+ 본 분석"),
        ("6", "표시", "홈 · 요약\n캘린더 · 알람"),
    ]
    for i, (n, h, b) in enumerate(steps):
        x = Inches(0.45) + i * Inches(2.15)
        if i == 4:
            add_rect(s, x, Inches(1.32), Inches(1.95), Inches(1.85), BLUE, 0.12)
        else:
            add_outline(s, x, Inches(1.32), Inches(1.95), Inches(1.85), WHITE, LINE, 0.12, 1.15)
        nc = BLUE if i != 4 else WHITE
        hc = TEXT if i != 4 else WHITE
        bc = MUTED if i != 4 else RGBColor(0xDB, 0xEA, 0xFE)
        add_text(s, x + Inches(0.12), Inches(1.42), Inches(1.7), Inches(0.28), n, 12, nc, True)
        add_text(s, x + Inches(0.12), Inches(1.72), Inches(1.7), Inches(0.35), h, 16, hc, True)
        add_text(s, x + Inches(0.12), Inches(2.12), Inches(1.7), Inches(0.85), b, 12, bc)
        if i < 5:
            arrow_right(s, x + Inches(1.92), Inches(2.1), Inches(0.26), Inches(0.16), BLUE)

    # detail split
    add_outline(s, Inches(0.45), Inches(3.45), Inches(6.05), Inches(3.05), WHITE, LINE, 0.08)
    add_text(s, Inches(0.7), Inches(3.62), Inches(5.6), Inches(0.32), "필터에 걸리면", 15, RED, True)
    add_lines(
        s,
        Inches(0.7),
        Inches(4.05),
        Inches(5.55),
        Inches(2.15),
        [
            "잡담 · 광고 · 로그인 알림",
            "본 분석을 건너뛰고 Room에서 삭제",
            "홈 목록에 남지 않음",
            "중요한 알림을 숨기지 않도록, 필터 실패 시에는 통과",
        ],
        14,
        MUTED,
    )

    add_outline(s, Inches(6.75), Inches(3.45), Inches(6.05), Inches(3.05), WHITE, LINE, 0.08)
    add_text(s, Inches(7.0), Inches(3.62), Inches(5.6), Inches(0.32), "분석에 통과하면", 15, GREEN, True)
    add_lines(
        s,
        Inches(7.0),
        Inches(4.05),
        Inches(5.55),
        Inches(2.15),
        [
            "요약 제목 · 유형 · 중요 여부 · 마감 시각",
            "유형: 수업 / 과제 / 소통 / 금전 / 기타",
            "마감이 있으면 캘린더 일정으로 연결",
            "사용자는 그 알림에 리마인더를 걸 수 있음",
        ],
        14,
        MUTED,
    )
    footer(s, 12)


def slide_flow_detail(prs):
    """A more visual sequence diagram style slide."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    kicker(s, "03  시스템 구조  ·  한 장 더")
    title(s, "누가 무엇을 하는가")

    cols = [
        ("타겟 앱", ["알림을 띄움"]),
        ("Android", ["Listener 수집", "Room 저장", "Worker 예약", "화면 표시"]),
        ("Backend", ["API 키 확인", "필터 모델 호출", "본 분석 호출", "결과만 반환"]),
        ("OpenAI", ["isRelevant?", "요약 JSON"]),
        ("사용자", ["홈에서 확인", "리마인더 설정", "캘린더 확인"]),
    ]
    colors = [NAVY2, BLUE, RGBColor(0x1D, 0x4E, 0xD8), RGBColor(0x10, 0xB9, 0x81), ORANGE]
    for i, ((name, items), c) in enumerate(zip(cols, colors)):
        x = Inches(0.45) + i * Inches(2.55)
        add_rect(s, x, Inches(1.32), Inches(2.38), Inches(0.55), c, 0.16)
        add_text(s, x, Inches(1.42), Inches(2.38), Inches(0.4), name, 15, WHITE, True, PP_ALIGN.CENTER)
        add_rect(s, x + Inches(1.1), Inches(1.87), Inches(0.04), Inches(4.55), LINE)
        for j, item in enumerate(items):
            y = Inches(2.15) + j * Inches(1.0)
            add_outline(s, x + Inches(0.12), y, Inches(2.14), Inches(0.78), WHITE, LINE, 0.16)
            add_text(s, x + Inches(0.2), y + Inches(0.22), Inches(1.98), Inches(0.4), item, 13, TEXT, False, PP_ALIGN.CENTER)

    footer(s, 13)


def slide_demo(prs, images: dict[str, Path]):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    content_bg(s)
    kicker(s, "04  시연")
    title(s, "이제 앱으로 이어집니다")

    add_text(
        s,
        Inches(0.6),
        Inches(1.15),
        Inches(12),
        Inches(0.4),
        "발표를 잠시 멈추고, 실제 화면을 보여드리겠습니다.",
        16,
        MUTED,
    )

    order = [
        ("홈", images["home"], "알림이 정리된 피드"),
        ("요약", images["summary"], "오늘 할 일 · Insight"),
        ("캘린더", images["calendar"], "마감이 찍힌 달력"),
        ("설정", images["keywords"], "키워드 · 타겟 앱"),
    ]
    for i, (name, img, cap) in enumerate(order):
        x = Inches(0.55) + i * Inches(3.2)
        add_text(s, x, Inches(1.6), Inches(2.9), Inches(0.32), f"{i+1}  {name}", 14, BLUE, True, PP_ALIGN.CENTER)
        phone(s, img, x + Inches(0.45), Inches(2.0), Inches(4.35))
        add_text(s, x, Inches(6.5), Inches(2.9), Inches(0.35), cap, 12, MUTED, False, PP_ALIGN.CENTER)

    footer(s, 14)


def slide_thanks(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(s, 0, 0, W, H, NAVY)
    add_rect(s, 0, 0, Inches(0.18), H, BLUE)
    blob = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.4), Inches(4.8), Inches(4.0), Inches(4.0))
    set_fill(blob, RGBColor(0x1E, 0x3A, 0x8A))
    blob2 = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.6), Inches(-1.4), Inches(4.2), Inches(4.2))
    set_fill(blob2, RGBColor(0x17, 0x24, 0x54))

    add_text(s, Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.1), "감사합니다", 52, WHITE, True)
    add_text(s, Inches(0.9), Inches(3.55), Inches(11.5), Inches(0.45), "질문 있으시면 편하게 해주세요.", 20, RGBColor(0x93, 0xC5, 0xFD))
    add_rect(s, Inches(0.9), Inches(4.35), Inches(1.8), Inches(0.06), BLUE, 0.5)
    add_text(s, Inches(0.9), Inches(4.6), Inches(10), Inches(0.4), "MyNoti  ·  Smarter Alerts, Zero Noise.", 14, RGBColor(0x94, 0xA3, 0xB8))


def main():
    rounded = {}
    mapping = {
        "home": ASSETS / "home.png",
        "summary": ASSETS / "summary.png",
        "calendar": ASSETS / "calendar.png",
        "apps": ASSETS / "settings_apps.png",
        "keywords": ASSETS / "settings_keywords.png",
    }
    for key, src in mapping.items():
        dst = ASSETS / f"{key}_round.png"
        round_phone(src, dst)
        rounded[key] = dst

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide_cover(prs)
    slide_agenda(prs)
    slide_morning(prs)
    slide_dilemma(prs)
    slide_problem(prs)
    slide_features(prs)
    slide_analyze(prs, rounded["home"])
    slide_summary(prs, rounded["summary"])
    slide_alarm(prs, rounded["keywords"], rounded["apps"])
    slide_calendar(prs, rounded["calendar"])
    slide_architecture(prs)
    slide_flow(prs)
    slide_flow_detail(prs)
    slide_demo(prs, rounded)
    slide_thanks(prs)

    prs.save(OUT)
    print(f"saved: {OUT}")


if __name__ == "__main__":
    main()
